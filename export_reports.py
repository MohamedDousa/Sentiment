#!/usr/bin/env python
"""
Export Reports for Staff Feedback Analysis Dashboard

This script provides enhanced PDF and PowerPoint export functionality for the dashboard.
"""

import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import seaborn as sns
import io
from datetime import datetime
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

# Import the processing functions
from preprocessing import load_data, preprocess_data, group_by_department
from nlp_pipeline import process_nlp_pipeline, NLPProcessor

# Output directory
OUTPUT_DIR = "reports"

def ensure_output_dir():
    """Ensure the output directory exists, create if it doesn't"""
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, "images"), exist_ok=True)

def fig_to_image(fig, filename, width=1000, height=600, format="png"):
    """Convert plotly figure to image file and return the path"""
    img_path = os.path.join(OUTPUT_DIR, "images", filename)
    fig.write_image(img_path, width=width, height=height, scale=2)
    return img_path

def generate_sentiment_chart(data):
    """Generate sentiment distribution chart"""
    # Create sentiment data
    sentiment_dist = {
        "positive": len(data[data['sentiment_score'] > 0.6]),
        "neutral": len(data[(data['sentiment_score'] >= 0.4) & (data['sentiment_score'] <= 0.6)]),
        "negative": len(data[data['sentiment_score'] < 0.4])
    }
    
    sentiment_data = []
    for sentiment, count in sentiment_dist.items():
        sentiment_data.append({
            "sentiment": sentiment.capitalize(),
            "count": count
        })
    
    sentiment_df = pd.DataFrame(sentiment_data)
    
    # Create Plotly pie chart
    fig = px.pie(
        sentiment_df,
        values="count",
        names="sentiment",
        title="Overall Sentiment Distribution",
        color="sentiment",
        color_discrete_map={
            "Positive": "green",
            "Neutral": "gray",
            "Negative": "red"
        }
    )
    fig.update_layout(
        title_font_size=22,
        font_size=16,
        legend_title_font_size=18
    )
    
    return fig, sentiment_df

def generate_themes_chart(data):
    """Generate theme overview chart"""
    # Create theme data for visualization
    theme_data = []
    for column in data.columns:
        if column.startswith('theme_') and column != 'theme_comment':
            theme_name = column.replace('theme_', '')
            # Skip themes containing "comment" as they skew results
            if "comment" not in theme_name.lower():
                count = data[column].sum()
                theme_data.append({
                    "theme": theme_name.replace("_", " ").title(),
                    "count": int(count)
                })
    
    theme_df = pd.DataFrame(theme_data)
    theme_df = theme_df.sort_values("count", ascending=False)
    
    # Display top 10 themes
    fig = px.bar(
        theme_df.head(10),
        x="theme",
        y="count",
        title="Top 10 Themes (excluding 'comment' themes)",
        labels={"theme": "Theme", "count": "Occurrence Count"},
        color="count",
        color_continuous_scale="Blues"
    )
    fig.update_layout(
        xaxis_tickangle=-45,
        title_font_size=22,
        font_size=16,
        xaxis_title_font_size=18,
        yaxis_title_font_size=18
    )
    
    return fig, theme_df

def generate_civility_themes_chart(data):
    """Generate civility and respect themes chart"""
    # Create theme data
    theme_data = []
    for column in data.columns:
        if column.startswith('theme_') and column != 'theme_comment':
            theme_name = column.replace('theme_', '')
            # Skip themes containing "comment" as they skew results
            if "comment" not in theme_name.lower():
                count = data[column].sum()
                theme_data.append({
                    "theme": theme_name.replace("_", " ").title(),
                    "count": int(count)
                })
    
    theme_df = pd.DataFrame(theme_data)
    
    # Filter for civility/respect related themes
    respect_themes = [
        'respect_communication', 'workplace_culture', 'leadership_behavior',
        'workplace_policies', 'inclusion_diversity', 'training_development',
        'recognition_appreciation'
    ]
    
    # Convert to display names
    respect_display_themes = [theme.replace("_", " ").title() for theme in respect_themes]
    
    respect_theme_df = theme_df[theme_df['theme'].isin(respect_display_themes)]
    respect_theme_df = respect_theme_df.sort_values("count", ascending=False)
    
    if not respect_theme_df.empty:
        fig = px.bar(
            respect_theme_df,
            x="theme",
            y="count",
            title="Civility & Respect Themes",
            labels={"theme": "Theme", "count": "Occurrence Count"},
            color="count",
            color_continuous_scale="Greens"
        )
        fig.update_layout(
            xaxis_tickangle=-45,
            title_font_size=22,
            font_size=16,
            xaxis_title_font_size=18,
            yaxis_title_font_size=18
        )
        
        return fig, respect_theme_df
    return None, None

def generate_theme_sentiment_chart(data):
    """Generate theme sentiment analysis chart"""
    # Get top 10 themes
    theme_data = []
    for column in data.columns:
        if column.startswith('theme_') and column != 'theme_comment':
            theme_name = column.replace('theme_', '')
            # Skip themes containing "comment" as they skew results
            if "comment" not in theme_name.lower():
                count = data[column].sum()
                theme_data.append({
                    "theme": theme_name,
                    "count": int(count)
                })
    
    theme_df = pd.DataFrame(theme_data)
    theme_df = theme_df.sort_values("count", ascending=False)
    top_themes = theme_df.head(8)['theme'].tolist()  # Limit to 8 for better readability
    
    # For each theme, analyze sentiment
    theme_sentiment_data = {}
    for theme in top_themes:
        theme_col = f"theme_{theme}"
        if theme_col in data.columns:
            # Get comments with this theme
            theme_comments = data[data[theme_col] == 1]
            
            # Count sentiment categories
            sentiment_counts = {
                "POSITIVE": len(theme_comments[theme_comments['sentiment_score'] > 0.6]),
                "NEUTRAL": len(theme_comments[(theme_comments['sentiment_score'] >= 0.4) & (theme_comments['sentiment_score'] <= 0.6)]),
                "NEGATIVE": len(theme_comments[theme_comments['sentiment_score'] < 0.4])
            }
            
            theme_sentiment_data[theme] = sentiment_counts
    
    # Create visualization data
    sentiment_rows = []
    for theme, counts in theme_sentiment_data.items():
        total = sum(counts.values())
        
        if total > 0:
            # Calculate percentages
            pos_pct = (counts["POSITIVE"] / total * 100) if total > 0 else 0
            neu_pct = (counts["NEUTRAL"] / total * 100) if total > 0 else 0
            neg_pct = (counts["NEGATIVE"] / total * 100) if total > 0 else 0
            
            # Format theme name for display
            display_theme = theme.replace("_", " ").title()
            
            # Add row for each sentiment category
            sentiment_rows.append({"theme": display_theme, "sentiment": "Positive", "percentage": pos_pct, "count": counts["POSITIVE"]})
            sentiment_rows.append({"theme": display_theme, "sentiment": "Neutral", "percentage": neu_pct, "count": counts["NEUTRAL"]})
            sentiment_rows.append({"theme": display_theme, "sentiment": "Negative", "percentage": neg_pct, "count": counts["NEGATIVE"]})
    
    # Create DataFrame
    sentiment_df = pd.DataFrame(sentiment_rows)
    
    # Create a grouped bar chart specifically for positive and negative percentages
    pivot_df = sentiment_df.pivot(index="theme", columns="sentiment", values="percentage").reset_index()
    
    # Sort by negative percentage (descending)
    if 'Negative' in pivot_df.columns:
        pivot_df = pivot_df.sort_values("Negative", ascending=False)
    
    fig = px.bar(
        pivot_df,
        x="theme",
        y=["Negative", "Neutral", "Positive"] if all(col in pivot_df.columns for col in ["Negative", "Neutral", "Positive"]) else pivot_df.columns[1:],
        title="Sentiment Distribution by Theme",
        labels={"theme": "Theme", "value": "Percentage (%)", "variable": "Sentiment"},
        color_discrete_map={
            "Positive": "green",
            "Neutral": "gray",
            "Negative": "red"
        },
        barmode="group"
    )
    
    fig.update_layout(
        xaxis_tickangle=-45,
        title_font_size=22,
        font_size=16,
        xaxis_title_font_size=18,
        yaxis_title_font_size=18,
        legend_title_font_size=18
    )
    
    return fig, sentiment_df

def generate_solutions_chart(data):
    """Generate solutions chart"""
    # Define solution categories with keywords
    specific_solutions = {
        "better_communication": {
            "keywords": ["better communication", "improve communication", "communicate better", "open communication"],
            "display_name": "Better Communication"
        },
        "management_training": {
            "keywords": ["management training", "train managers", "leadership training", "leader training"],
            "display_name": "Management/Leadership Training"
        },
        "staff_training": {
            "keywords": ["staff training", "employee training", "training for staff", "train employees"],
            "display_name": "Staff Training & Development"
        },
        "accountability": {
            "keywords": ["accountability", "hold accountable", "consequences", "responsible", "responsibility"],
            "display_name": "Accountability Measures"
        },
        "culture_improvement": {
            "keywords": ["better culture", "improve culture", "positive culture", "workplace culture"],
            "display_name": "Culture Improvement"
        },
        "lead_by_example": {
            "keywords": ["lead by example", "leading by example", "role model", "role models", "role modeling"],
            "display_name": "Leading by Example"
        },
        "teambuilding": {
            "keywords": ["team building", "team activities", "social events", "get together", "social activities"],
            "display_name": "Team Building Activities"
        },
        "recognition_programs": {
            "keywords": ["recognition", "rewards", "incentives", "award", "appreciate", "appreciation"],
            "display_name": "Recognition & Appreciation"
        }
    }
    
    # Count comments mentioning each solution
    for solution_key, solution_data in specific_solutions.items():
        comment_count = 0
        for keyword in solution_data["keywords"]:
            # Check for keyword in free_text_comments
            matches = data[data['free_text_comments'].str.contains(keyword, case=False, na=False)]
            comment_count += len(matches)
        
        specific_solutions[solution_key]["count"] = comment_count
    
    # Create visualization data
    solution_counts = []
    for sol_key, sol_data in specific_solutions.items():
        count = sol_data["count"]
        if count > 0:
            solution_counts.append({
                "solution": sol_data["display_name"],
                "count": count
            })
    
    # Sort solutions by count and take top 6
    solution_df = pd.DataFrame(solution_counts)
    solution_df = solution_df.sort_values("count", ascending=False)
    solution_df = solution_df.head(6)  # Focus on top 6 solutions for better readability
    
    # Calculate percentage of total comments
    total_comments = len(data)
    solution_df["percentage"] = (solution_df["count"] / total_comments * 100).round(1)
    
    # Create a bar chart
    fig = px.bar(
        solution_df,
        x="solution",
        y="count",
        title="Top Solutions Mentioned in Feedback",
        labels={"solution": "Solution Type", "count": "Number of Comments"},
        color="count",
        color_continuous_scale="Viridis",
        text="percentage"
    )
    fig.update_traces(texttemplate='%{text}%', textposition='outside')
    fig.update_layout(
        xaxis_tickangle=-45,
        title_font_size=22,
        font_size=16,
        xaxis_title_font_size=18,
        yaxis_title_font_size=18
    )
    
    return fig, solution_df

def generate_pdf_report(data_file, output_file=None):
    """Generate PDF report with all visualizations"""
    print(f"Generating PDF report from data file: {data_file}")
    
    # Create output directory
    ensure_output_dir()
    
    if not output_file:
        output_file = os.path.join(OUTPUT_DIR, f"feedback_analysis_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
    
    try:
        # Load and process data
        print("Loading and processing data...")
        raw_data = load_data(data_file)
        processed_data = preprocess_data(raw_data)
        
        # Run NLP pipeline
        print("Running NLP pipeline...")
        nlp_processor = NLPProcessor()
        processed_data = process_nlp_pipeline(processed_data, nlp_processor)
        
        # Create PDF object
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        
        # Add title
        pdf.set_font("Arial", "B", 24)
        pdf.cell(0, 20, "Staff Feedback Analysis Report", ln=True, align="C")
        
        # Add generation date
        pdf.set_font("Arial", "", 12)
        pdf.cell(0, 10, f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", ln=True, align="C")
        
        # Add summary information
        pdf.ln(10)
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "Data Summary", ln=True)
        
        pdf.set_font("Arial", "", 12)
        pdf.cell(0, 10, f"Total comments analyzed: {len(processed_data)}", ln=True)
        
        # Key metrics
        sentiment_dist = {
            "positive": len(processed_data[processed_data['sentiment_score'] > 0.6]),
            "neutral": len(processed_data[(processed_data['sentiment_score'] >= 0.4) & (processed_data['sentiment_score'] <= 0.6)]),
            "negative": len(processed_data[processed_data['sentiment_score'] < 0.4])
        }
        
        pdf.cell(0, 10, f"Positive comments: {sentiment_dist['positive']} ({sentiment_dist['positive']/len(processed_data)*100:.1f}%)", ln=True)
        pdf.cell(0, 10, f"Neutral comments: {sentiment_dist['neutral']} ({sentiment_dist['neutral']/len(processed_data)*100:.1f}%)", ln=True)
        pdf.cell(0, 10, f"Negative comments: {sentiment_dist['negative']} ({sentiment_dist['negative']/len(processed_data)*100:.1f}%)", ln=True)
        
        # Generate and add charts
        print("Generating charts for PDF...")
        
        # Sentiment distribution chart
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "1. Sentiment Distribution", ln=True)
        
        sentiment_fig, _ = generate_sentiment_chart(processed_data)
        sentiment_img = fig_to_image(sentiment_fig, "sentiment_distribution.png")
        pdf.image(sentiment_img, x=10, y=30, w=190)
        
        # Themes overview chart
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "2. Top Themes", ln=True)
        
        themes_fig, themes_df = generate_themes_chart(processed_data)
        themes_img = fig_to_image(themes_fig, "themes_overview.png")
        pdf.image(themes_img, x=10, y=30, w=190)
        
        # Add theme table
        pdf.ln(180)  # Move down after the image
        pdf.set_font("Arial", "B", 12)
        pdf.cell(95, 10, "Theme", 1, 0, "C")
        pdf.cell(95, 10, "Count", 1, 1, "C")
        
        pdf.set_font("Arial", "", 12)
        for idx, row in themes_df.head(10).iterrows():
            pdf.cell(95, 10, row["theme"], 1, 0)
            pdf.cell(95, 10, str(row["count"]), 1, 1, "C")
        
        # Civility themes chart
        civility_fig, civility_df = generate_civility_themes_chart(processed_data)
        if civility_fig is not None:
            pdf.add_page()
            pdf.set_font("Arial", "B", 16)
            pdf.cell(0, 10, "3. Civility & Respect Themes", ln=True)
            
            civility_img = fig_to_image(civility_fig, "civility_themes.png")
            pdf.image(civility_img, x=10, y=30, w=190)
        
        # Theme sentiment analysis
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "4. Theme Sentiment Analysis", ln=True)
        
        theme_sentiment_fig, _ = generate_theme_sentiment_chart(processed_data)
        theme_sentiment_img = fig_to_image(theme_sentiment_fig, "theme_sentiment.png")
        pdf.image(theme_sentiment_img, x=10, y=30, w=190)
        
        # Solutions
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "5. Key Solutions", ln=True)
        
        solutions_fig, solutions_df = generate_solutions_chart(processed_data)
        solutions_img = fig_to_image(solutions_fig, "solutions.png")
        pdf.image(solutions_img, x=10, y=30, w=190)
        
        # Add solutions table
        pdf.ln(180)  # Move down after the image
        pdf.set_font("Arial", "B", 12)
        pdf.cell(65, 10, "Solution", 1, 0, "C")
        pdf.cell(65, 10, "Count", 1, 0, "C")
        pdf.cell(60, 10, "Percentage", 1, 1, "C")
        
        pdf.set_font("Arial", "", 12)
        for idx, row in solutions_df.iterrows():
            pdf.cell(65, 10, row["solution"], 1, 0)
            pdf.cell(65, 10, str(row["count"]), 1, 0, "C")
            pdf.cell(60, 10, f"{row['percentage']}%", 1, 1, "C")
        
        # Actionable insights
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "6. Actionable Insights", ln=True)
        
        pdf.ln(10)
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "Recommended Actions Based on Feedback:", ln=True)
        
        # Get top solutions
        top_solutions = solutions_df.head(3)["solution"].tolist()
        
        pdf.set_font("Arial", "", 12)
        for i, solution in enumerate(top_solutions):
            pdf.ln(5)
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 10, f"{i+1}. {solution}", ln=True)
            
            pdf.set_font("Arial", "", 12)
            if "Communication" in solution:
                pdf.multi_cell(0, 10, "- Implement regular team meetings focused on open dialogue\n- Create clear communication channels for feedback\n- Provide training on effective communication strategies")
            elif "Training" in solution:
                pdf.multi_cell(0, 10, "- Develop targeted training programs for managers and staff\n- Incorporate respect and civility into existing training\n- Provide ongoing education about workplace expectations")
            elif "Example" in solution:
                pdf.multi_cell(0, 10, "- Ensure managers model respectful behavior\n- Recognize and highlight positive examples\n- Create accountability for leadership behaviors")
            elif "Recognition" in solution:
                pdf.multi_cell(0, 10, "- Create a formal recognition program\n- Implement peer-to-peer appreciation systems\n- Regularly acknowledge positive behaviors")
            elif "Culture" in solution:
                pdf.multi_cell(0, 10, "- Conduct a culture assessment\n- Define clear values around respect\n- Involve staff in creating cultural norms")
            elif "Accountability" in solution:
                pdf.multi_cell(0, 10, "- Establish clear expectations for behavior\n- Create consistent response to incivility\n- Implement fair enforcement of policies")
        
        # Save the PDF
        pdf.output(output_file)
        print(f"PDF report generated and saved to {output_file}")
        return output_file
    
    except Exception as e:
        print(f"Error generating PDF report: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

def generate_pptx_report(data_file, output_file=None):
    """Generate PowerPoint report with all visualizations"""
    print(f"Generating PowerPoint report from data file: {data_file}")
    
    # Create output directory
    ensure_output_dir()
    
    if not output_file:
        output_file = os.path.join(OUTPUT_DIR, f"feedback_analysis_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    
    try:
        # Load and process data
        print("Loading and processing data...")
        raw_data = load_data(data_file)
        processed_data = preprocess_data(raw_data)
        
        # Run NLP pipeline
        print("Running NLP pipeline...")
        nlp_processor = NLPProcessor()
        processed_data = process_nlp_pipeline(processed_data, nlp_processor)
        
        # Create presentation
        prs = Presentation()
        
        # SLIDE 1: Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Staff Feedback Analysis Report"
        subtitle.text = f"Generated on: {datetime.now().strftime('%Y-%m-%d')}"
        
        # SLIDE 2: Introduction
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = "Report Overview"
        tf = body.text_frame
        tf.text = "This report provides analysis of staff feedback on workplace civility and respect."
        
        p = tf.add_paragraph()
        p.text = f"Total comments analyzed: {len(processed_data)}"
        p.level = 1
        
        # Add key metrics
        sentiment_dist = {
            "positive": len(processed_data[processed_data['sentiment_score'] > 0.6]),
            "neutral": len(processed_data[(processed_data['sentiment_score'] >= 0.4) & (processed_data['sentiment_score'] <= 0.6)]),
            "negative": len(processed_data[processed_data['sentiment_score'] < 0.4])
        }
        
        p = tf.add_paragraph()
        p.text = f"Sentiment breakdown:"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = f"Positive: {sentiment_dist['positive']} ({sentiment_dist['positive']/len(processed_data)*100:.1f}%)"
        p.level = 2
        
        p = tf.add_paragraph()
        p.text = f"Neutral: {sentiment_dist['neutral']} ({sentiment_dist['neutral']/len(processed_data)*100:.1f}%)"
        p.level = 2
        
        p = tf.add_paragraph()
        p.text = f"Negative: {sentiment_dist['negative']} ({sentiment_dist['negative']/len(processed_data)*100:.1f}%)"
        p.level = 2
        
        # Generate and add charts
        print("Generating charts for PowerPoint...")
        
        # SLIDE 3: Sentiment Distribution
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout with title and content
        title = slide.shapes.title
        title.text = "Sentiment Distribution"
        
        sentiment_fig, _ = generate_sentiment_chart(processed_data)
        sentiment_img = fig_to_image(sentiment_fig, "sentiment_distribution.png")
        
        # Add image to slide
        slide.shapes.add_picture(sentiment_img, Inches(1), Inches(1.5), width=Inches(8))
        
        # SLIDE 4: Top Themes
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Top Themes from Feedback"
        
        themes_fig, themes_df = generate_themes_chart(processed_data)
        themes_img = fig_to_image(themes_fig, "themes_overview.png")
        
        # Add image to slide
        slide.shapes.add_picture(themes_img, Inches(1), Inches(1.5), width=Inches(8))
        
        # SLIDE 5: Civility Themes
        civility_fig, civility_df = generate_civility_themes_chart(processed_data)
        if civility_fig is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Civility & Respect Themes"
            
            civility_img = fig_to_image(civility_fig, "civility_themes.png")
            slide.shapes.add_picture(civility_img, Inches(1), Inches(1.5), width=Inches(8))
        
        # SLIDE 6: Theme Sentiment Analysis
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Theme Sentiment Analysis"
        
        theme_sentiment_fig, _ = generate_theme_sentiment_chart(processed_data)
        theme_sentiment_img = fig_to_image(theme_sentiment_fig, "theme_sentiment.png")
        slide.shapes.add_picture(theme_sentiment_img, Inches(1), Inches(1.5), width=Inches(8))
        
        # SLIDE 7: Key Solutions
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Key Solutions from Feedback"
        
        solutions_fig, solutions_df = generate_solutions_chart(processed_data)
        solutions_img = fig_to_image(solutions_fig, "solutions.png")
        slide.shapes.add_picture(solutions_img, Inches(1), Inches(1.5), width=Inches(8))
        
        # SLIDE 8-10: Actionable Insights - one slide per top solution
        top_solutions = solutions_df.head(3)["solution"].tolist()
        
        for i, solution in enumerate(top_solutions):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content with bullets
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = f"Actionable Insight: {solution}"
            tf = body.text_frame
            
            if "Communication" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Implement regular team meetings focused on open dialogue"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Create clear communication channels for feedback"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Provide training on effective communication strategies"
                p.level = 1
                
            elif "Training" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Develop targeted training programs for managers and staff"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Incorporate respect and civility into existing training"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Provide ongoing education about workplace expectations"
                p.level = 1
                
            elif "Example" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Ensure managers model respectful behavior"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Recognize and highlight positive examples"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Create accountability for leadership behaviors"
                p.level = 1
                
            elif "Recognition" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Create a formal recognition program"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Implement peer-to-peer appreciation systems"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Regularly acknowledge positive behaviors"
                p.level = 1
                
            elif "Culture" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Conduct a culture assessment"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Define clear values around respect"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Involve staff in creating cultural norms"
                p.level = 1
                
            elif "Accountability" in solution:
                tf.text = "Recommended actions:"
                
                p = tf.add_paragraph()
                p.text = "Establish clear expectations for behavior"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Create consistent response to incivility"
                p.level = 1
                
                p = tf.add_paragraph()
                p.text = "Implement fair enforcement of policies"
                p.level = 1
        
        # SLIDE 11: Summary and Next Steps
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = "Summary & Next Steps"
        tf = body.text_frame
        tf.text = "Key Takeaways:"
        
        p = tf.add_paragraph()
        p.text = f"Focus on implementing {', '.join(top_solutions[:-1])} and {top_solutions[-1]}"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Prioritize specific themes identified in the analysis"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Conduct follow-up surveys to measure progress"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Next Steps:"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "Create action plan based on key insights"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Assign team members to specific initiatives"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Set measurable goals and timelines"
        p.level = 1
        
        # Save the presentation
        prs.save(output_file)
        print(f"PowerPoint report generated and saved to {output_file}")
        return output_file
    
    except Exception as e:
        print(f"Error generating PowerPoint report: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Generate reports from feedback analysis")
    parser.add_argument("--data", "-d", required=True, help="Path to data file (CSV or Excel)")
    parser.add_argument("--output-dir", "-o", help="Output directory (default: reports)")
    parser.add_argument("--format", "-f", choices=["pdf", "pptx", "both"], default="both", help="Report format to generate")
    
    args = parser.parse_args()
    
    if args.output_dir:
        OUTPUT_DIR = args.output_dir
    
    # Run report generation
    if args.format in ["pdf", "both"]:
        pdf_file = generate_pdf_report(args.data)
        if pdf_file:
            print(f"PDF report saved to: {pdf_file}")
    
    if args.format in ["pptx", "both"]:
        pptx_file = generate_pptx_report(args.data)
        if pptx_file:
            print(f"PowerPoint report saved to: {pptx_file}") 